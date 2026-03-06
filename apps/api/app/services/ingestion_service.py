"""
Document ingestion service.

Pipeline:
  File → Text Extraction → Chunking → Embedding → Qdrant Storage → DB Metadata
"""
import uuid
from datetime import datetime, timezone
from pathlib import Path

from qdrant_client import AsyncQdrantClient
from qdrant_client.models import (
    Distance,
    PointStruct,
    SparseIndexParams,
    SparseVectorParams,
    VectorParams,
    VectorsConfig,
)
from sqlalchemy.ext.asyncio import AsyncSession

from app.core.config import settings
from app.db.models.document import Document
from app.db.models.document_chunk import DocumentChunk
from app.db.models.knowledge_base import KnowledgeBase
from app.services.embedding_service import EmbeddingService


MAX_CHUNK_SIZE = settings.CHUNK_SIZE
CHUNK_OVERLAP = settings.CHUNK_OVERLAP


class IngestionService:
    """
    Orchestrates the full document ingestion pipeline.
    Each step is isolated and can be run independently.
    """

    def __init__(
        self,
        db: AsyncSession,
        qdrant_client: AsyncQdrantClient,
        embedding_service: EmbeddingService,
    ):
        self._db = db
        self._qdrant = qdrant_client
        self._embedding_service = embedding_service

    async def ingest_document(
        self,
        document_id: uuid.UUID,
    ) -> None:
        """Full ingestion pipeline for a single document."""
        doc = await self._db.get(Document, document_id)
        if not doc:
            raise ValueError(f"Document {document_id} not found")

        kb = await self._db.get(KnowledgeBase, doc.kb_id)
        if not kb:
            raise ValueError(f"KnowledgeBase {doc.kb_id} not found")

        # Update status to processing
        doc.status = "processing"
        await self._db.commit()

        try:
            # 1. Text extraction
            text_pages = extract_text(doc.file_path, doc.file_type)
            doc.page_count = len(text_pages)

            # 2. Chunking
            chunks = chunk_pages(text_pages, MAX_CHUNK_SIZE, CHUNK_OVERLAP)

            # 3. Ensure Qdrant collection exists
            await self._ensure_collection(kb.qdrant_collection)

            # 4. Embed + store in batches
            chunk_records = []
            points = []

            batch_size = settings.EMBEDDING_BATCH_SIZE
            for batch_start in range(0, len(chunks), batch_size):
                batch = chunks[batch_start: batch_start + batch_size]
                texts = [c["content"] for c in batch]
                dense_vectors = await self._embedding_service.embed_documents(texts)

                for chunk_data, dense_vec in zip(batch, dense_vectors):
                    point_id = str(uuid.uuid4())
                    sparse_indices, sparse_values = (
                        self._embedding_service.compute_sparse_vector(chunk_data["content"])
                    )

                    # Qdrant point
                    points.append(
                        PointStruct(
                            id=point_id,
                            vector={
                                "dense": dense_vec,
                                "sparse": {
                                    "indices": sparse_indices,
                                    "values": sparse_values,
                                },
                            },
                            payload={
                                "content": chunk_data["content"],
                                "document_id": str(doc.id),
                                "document_title": doc.filename,
                                "kb_id": str(kb.id),
                                "kb_name": kb.name,
                                "page_number": chunk_data.get("page_number"),
                                "chunk_index": chunk_data.get("chunk_index", 0),
                                "source_url": chunk_data.get("source_url"),
                            },
                        )
                    )

                    # DB chunk record
                    chunk_records.append(
                        DocumentChunk(
                            doc_id=doc.id,
                            kb_id=kb.id,
                            content=chunk_data["content"],
                            qdrant_point_id=point_id,
                            page_number=chunk_data.get("page_number"),
                            chunk_index=chunk_data.get("chunk_index", 0),
                            token_count=len(chunk_data["content"].split()),
                        )
                    )

            # 5. Upsert to Qdrant
            await self._qdrant.upsert(
                collection_name=kb.qdrant_collection,
                points=points,
            )

            # 6. Save chunk metadata to PostgreSQL
            self._db.add_all(chunk_records)
            doc.chunk_count = len(chunk_records)
            doc.status = "completed"
            doc.processed_at = datetime.now(timezone.utc)

            # Update KB aggregate counts
            kb.document_count = kb.document_count + 1
            kb.chunk_count = kb.chunk_count + len(chunk_records)

            await self._db.commit()

        except Exception as exc:
            doc.status = "failed"
            doc.error_message = str(exc)[:1000]
            await self._db.commit()
            raise

    async def _ensure_collection(self, collection_name: str) -> None:
        existing = await self._qdrant.collection_exists(collection_name)
        if existing:
            return

        await self._qdrant.create_collection(
            collection_name=collection_name,
            vectors_config={
                "dense": VectorParams(
                    size=settings.QDRANT_VECTOR_SIZE,
                    distance=Distance.COSINE,
                )
            },
            sparse_vectors_config={
                "sparse": SparseVectorParams(
                    index=SparseIndexParams(on_disk=False)
                )
            },
        )


# ---------------------------------------------------------------------------
# Text extraction
# ---------------------------------------------------------------------------

def extract_text(file_path: str, file_type: str) -> list[dict]:
    """
    Returns list of {page_number: int, content: str}.
    Supports: pdf, docx, pptx, txt, csv.
    """
    file_type = file_type.lower().lstrip(".")
    path = Path(file_path)

    if file_type == "pdf":
        return _extract_pdf(path)
    elif file_type == "docx":
        return _extract_docx(path)
    elif file_type == "pptx":
        return _extract_pptx(path)
    elif file_type in ("txt", "md"):
        return _extract_text_file(path)
    elif file_type == "csv":
        return _extract_csv(path)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")


def _extract_pdf(path: Path) -> list[dict]:
    import fitz  # PyMuPDF
    pages = []
    with fitz.open(str(path)) as doc:
        for page_num, page in enumerate(doc, start=1):
            text = page.get_text()
            if text.strip():
                pages.append({"page_number": page_num, "content": text})
    return pages


def _extract_docx(path: Path) -> list[dict]:
    from docx import Document as DocxDocument
    doc = DocxDocument(str(path))
    full_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return [{"page_number": 1, "content": full_text}]


def _extract_pptx(path: Path) -> list[dict]:
    from pptx import Presentation
    prs = Presentation(str(path))
    slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)
        if texts:
            slides.append({"page_number": idx, "content": "\n".join(texts)})
    return slides


def _extract_text_file(path: Path) -> list[dict]:
    content = path.read_text(encoding="utf-8", errors="replace")
    return [{"page_number": 1, "content": content}]


def _extract_csv(path: Path) -> list[dict]:
    import csv
    rows = []
    with open(path, newline="", encoding="utf-8") as f:
        reader = csv.reader(f)
        rows = [", ".join(row) for row in reader]
    return [{"page_number": 1, "content": "\n".join(rows)}]


# ---------------------------------------------------------------------------
# Chunking
# ---------------------------------------------------------------------------

def chunk_pages(
    pages: list[dict],
    chunk_size: int = MAX_CHUNK_SIZE,
    overlap: int = CHUNK_OVERLAP,
) -> list[dict]:
    """Split pages into overlapping token-based chunks."""
    chunks = []
    for page in pages:
        page_chunks = _split_text(
            page["content"],
            chunk_size=chunk_size,
            overlap=overlap,
            page_number=page["page_number"],
        )
        chunks.extend(page_chunks)
    return chunks


def _split_text(
    text: str, chunk_size: int, overlap: int, page_number: int
) -> list[dict]:
    words = text.split()
    result = []
    step = max(1, chunk_size - overlap)

    for i, start in enumerate(range(0, len(words), step)):
        chunk_words = words[start: start + chunk_size]
        if not chunk_words:
            break
        result.append(
            {
                "content": " ".join(chunk_words),
                "page_number": page_number,
                "chunk_index": i,
            }
        )

    return result
